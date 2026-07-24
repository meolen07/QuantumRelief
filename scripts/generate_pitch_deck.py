#!/usr/bin/env python3
"""Generate the QuantumRelief Quantrio pitch deck (10 slides, 16:9)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Design tokens
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x0B, 0x1C, 0x2C)
NAVY_MID = RGBColor(0x12, 0x2B, 0x42)
NAVY_LIGHT = RGBColor(0x1A, 0x3A, 0x55)
ORANGE = RGBColor(0xFF, 0x6B, 0x1A)
ORANGE_SOFT = RGBColor(0xFF, 0x8C, 0x42)
SLATE = RGBColor(0x94, 0xA3, 0xB8)
SLATE_LIGHT = RGBColor(0xCB, 0xD5, 0xE1)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
MUTED = RGBColor(0x64, 0x74, 0x8B)
CARD = RGBColor(0x14, 0x2E, 0x45)
CARD_BORDER = RGBColor(0x1E, 0x40, 0x5A)
ACCENT_TEAL = RGBColor(0x38, 0xBD, 0xC2)

MARGIN = Inches(0.55)
CONTENT_TOP = Inches(1.35)

ROOT = Path(__file__).resolve().parents[1]
OUT_PATH = ROOT / "QuantumRelief_Quantrio_Pitch.pptx"


def _set_run(run, *, size=18, bold=False, color=WHITE, font="Calibri"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_text(
    tf,
    text,
    *,
    size=18,
    bold=False,
    color=WHITE,
    font="Calibri",
    align=PP_ALIGN.LEFT,
    space_after=6,
    space_before=0,
    clear=True,
):
    if clear:
        tf.clear()
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    _set_run(run, size=size, bold=bold, color=color, font=font)
    return p


def _fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _line(shape, color: RGBColor, width_pt: float = 1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def _rect(slide, left, top, width, height, fill: RGBColor, line: RGBColor | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _fill(shape, fill)
    if line is not None:
        _line(shape, line, 1.0)
    else:
        shape.line.fill.background()
    return shape


def _round_rect(slide, left, top, width, height, fill: RGBColor, line: RGBColor | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    _fill(shape, fill)
    # Soften corners
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    if line is not None:
        _line(shape, line, 1.0)
    else:
        shape.line.fill.background()
    return shape


def _textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def _accent_bar(slide, left, top, width=Inches(0.9), height=Inches(0.06)):
    return _rect(slide, left, top, width, height, ORANGE)


def _dark_bg(slide):
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)


def _abstract_flow(slide, *, intensity: str = "medium"):
    """Draw abstract navy flow lines / grid map using shapes only."""
    # Soft large orbs
    for left, top, w, h, c in [
        (Inches(-1.2), Inches(-1.0), Inches(4.5), Inches(4.5), NAVY_MID),
        (Inches(10.2), Inches(4.2), Inches(4.8), Inches(4.8), NAVY_LIGHT),
        (Inches(8.8), Inches(-1.5), Inches(3.5), Inches(3.5), NAVY_MID),
    ]:
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, w, h)
        _fill(oval, c)

    # Diagonal accent strokes
    strokes = [
        (Inches(0.2), Inches(6.6), Inches(4.8), Inches(0.03)),
        (Inches(0.5), Inches(6.85), Inches(3.2), Inches(0.02)),
        (Inches(9.2), Inches(0.35), Inches(3.6), Inches(0.025)),
    ]
    for left, top, w, h in strokes:
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
        _fill(s, ORANGE if intensity != "subtle" else RGBColor(0xC4, 0x55, 0x12))
        s.rotation = -12.0

    # Subtle grid map (right side)
    if intensity != "subtle":
        base_l, base_t = Inches(9.6), Inches(1.8)
        for i in range(8):
            for j in range(6):
                if (i + j) % 3 == 0:
                    continue
                cell = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    base_l + Inches(i * 0.38),
                    base_t + Inches(j * 0.38),
                    Inches(0.28),
                    Inches(0.28),
                )
                _fill(cell, NAVY_LIGHT if (i + j) % 2 == 0 else CARD)
                cell.line.color.rgb = CARD_BORDER
                cell.line.width = Pt(0.5)

        # Flow path across grid
        path_pts = [(0, 5), (1, 4), (2, 4), (3, 3), (4, 2), (5, 2), (6, 1), (7, 0)]
        for idx, (i, j) in enumerate(path_pts):
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                base_l + Inches(i * 0.38 + 0.06),
                base_t + Inches(j * 0.38 + 0.06),
                Inches(0.16),
                Inches(0.16),
            )
            _fill(dot, ORANGE if idx == len(path_pts) - 1 else ORANGE_SOFT)


def _header(slide, eyebrow: str, title: str, *, subtitle: str | None = None):
    _accent_bar(slide, MARGIN, Inches(0.42))
    box = _textbox(slide, MARGIN, Inches(0.52), Inches(11.5), Inches(0.35))
    _add_text(box.text_frame, eyebrow.upper(), size=11, bold=True, color=ORANGE, space_after=0)

    tbox = _textbox(slide, MARGIN, Inches(0.82), Inches(12.0), Inches(0.55))
    _add_text(tbox.text_frame, title, size=28, bold=True, color=WHITE, space_after=0)

    if subtitle:
        sbox = _textbox(slide, MARGIN, Inches(1.32), Inches(12.0), Inches(0.35))
        _add_text(sbox.text_frame, subtitle, size=14, color=SLATE, space_after=0)


def _footer(slide, page: int, total: int = 10):
    line = _rect(slide, MARGIN, Inches(7.15), Inches(12.2), Inches(0.01), CARD_BORDER)
    left = _textbox(slide, MARGIN, Inches(7.2), Inches(6), Inches(0.28))
    _add_text(left.text_frame, "QuantumRelief  ·  Team 5 — Quantrio", size=10, color=MUTED, space_after=0)
    right = _textbox(slide, Inches(10.5), Inches(7.2), Inches(2.3), Inches(0.28))
    _add_text(
        right.text_frame,
        f"{page:02d} / {total:02d}",
        size=10,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
        space_after=0,
    )
    return line


def _bullet_block(tf, items: list[str], *, size=15, color=SLATE_LIGHT, bullet_color=ORANGE):
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        p.level = 0
        # Custom orange dash via run prefix
        run = p.add_run()
        run.text = "▸  "
        _set_run(run, size=size, bold=True, color=bullet_color)
        run2 = p.add_run()
        run2.text = item
        _set_run(run2, size=size, color=color)


def _style_chart(chart, *, value_axis_max=100):
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(11)
    chart.legend.font.color.rgb = SLATE_LIGHT

    plot = chart.plots[0]
    plot.gap_width = 80

    # Series colors
    colors = [ORANGE, SLATE]
    for i, series in enumerate(chart.series):
        try:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = colors[i % len(colors)]
        except Exception:
            pass

    # Value axis
    try:
        va = chart.value_axis
        va.has_major_gridlines = True
        va.major_gridlines.format.line.color.rgb = CARD_BORDER
        va.maximum_scale = value_axis_max
        va.minimum_scale = 80
        va.tick_labels.font.size = Pt(10)
        va.tick_labels.font.color.rgb = SLATE
        va.format.line.color.rgb = CARD_BORDER
    except Exception:
        pass

    try:
        ca = chart.category_axis
        ca.tick_labels.font.size = Pt(12)
        ca.tick_labels.font.color.rgb = SLATE_LIGHT
        ca.format.line.color.rgb = CARD_BORDER
    except Exception:
        pass


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------
def slide_01_title(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _abstract_flow(slide, intensity="strong")

    # Left brand panel accent
    _rect(slide, 0, 0, Inches(0.12), SLIDE_H, ORANGE)

    # Team pill
    pill = _round_rect(slide, MARGIN, Inches(1.55), Inches(5.4), Inches(0.42), CARD, CARD_BORDER)
    ptf = pill.text_frame
    ptf.word_wrap = True
    ptf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = ptf.paragraphs[0].add_run()
    run.text = "TEAM 5 — QUANTRIO  |  QC4SG — SEA HACKATHON"
    _set_run(run, size=11, bold=True, color=ORANGE)

    title = _textbox(slide, MARGIN, Inches(2.2), Inches(10), Inches(1.1))
    _add_text(title.text_frame, "QuantumRelief", size=54, bold=True, color=WHITE, space_after=0)

    _accent_bar(slide, MARGIN, Inches(3.35), width=Inches(1.4), height=Inches(0.07))

    sub = _textbox(slide, MARGIN, Inches(3.55), Inches(9.5), Inches(1.1))
    _add_text(
        sub.text_frame,
        "Navigating Disaster Chaos with\nQuantum-Classical Hybrid Intelligence.",
        size=22,
        color=SLATE_LIGHT,
        space_after=0,
    )

    foot = _textbox(slide, MARGIN, Inches(6.55), Inches(8), Inches(0.4))
    _add_text(
        foot.text_frame,
        "Emergency Escape Routing  ·  Manila / Intramuros Prototype",
        size=12,
        color=MUTED,
        space_after=0,
    )


def slide_02_problem(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _abstract_flow(slide, intensity="subtle")
    _header(
        slide,
        "The Problem",
        "The \"Big One\" is Inevitable",
        subtitle="Why current systems fail during high-magnitude earthquakes in Manila",
    )

    # Highlighted numbers
    stats = [
        ("7.2M", "People at risk\nin Metro Manila"),
        ("Seconds", "Static maps become\nuseless after impact"),
        ("Too Slow", "Classical routing lags\nbehind real-time chaos"),
    ]
    card_w = Inches(3.7)
    gap = Inches(0.35)
    start_x = MARGIN
    top = Inches(2.0)
    for i, (num, label) in enumerate(stats):
        left = start_x + i * (card_w + gap)
        card = _round_rect(slide, left, top, card_w, Inches(2.35), CARD, CARD_BORDER)
        # Top accent strip
        _rect(slide, left, top, card_w, Inches(0.08), ORANGE)

        nbox = _textbox(slide, left + Inches(0.25), top + Inches(0.45), card_w - Inches(0.5), Inches(0.8))
        _add_text(nbox.text_frame, num, size=40, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, space_after=0)

        lbox = _textbox(slide, left + Inches(0.3), top + Inches(1.35), card_w - Inches(0.6), Inches(0.85))
        tf = lbox.text_frame
        tf.word_wrap = True
        _add_text(tf, label, size=14, color=SLATE_LIGHT, align=PP_ALIGN.CENTER, space_after=0)

    # Bottom insight
    insight = _round_rect(slide, MARGIN, Inches(4.7), Inches(12.2), Inches(1.85), CARD, CARD_BORDER)
    _rect(slide, MARGIN, Inches(4.7), Inches(0.1), Inches(1.85), ORANGE)
    ibox = _textbox(slide, MARGIN + Inches(0.35), Inches(4.9), Inches(11.5), Inches(1.5))
    tf = ibox.text_frame
    tf.word_wrap = True
    _add_text(tf, "THE FAILURE MODE", size=11, bold=True, color=ORANGE, space_after=8)
    _add_text(
        tf,
        "Infrastructure maps assume a fixed world. In a major quake, roads collapse, "
        "bridges fail, and traffic surges toward exits — turning yesterday's optimal "
        "path into a deadly dead-end. First responders need decisions that anticipate "
        "chaos, not react to it minutes too late.",
        size=15,
        color=SLATE_LIGHT,
        space_after=0,
        clear=False,
    )
    _footer(slide, 2)


def slide_03_case_study(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)

    # Bleed right panel
    _rect(slide, Inches(7.4), 0, Inches(5.933), SLIDE_H, NAVY_MID)

    # Abstract "map chaos" on right
    for i in range(12):
        for j in range(14):
            if (i * 3 + j * 5) % 7 == 0:
                continue
            blocked = (i + j) % 4 == 0
            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(7.7) + Inches(i * 0.42),
                Inches(0.9) + Inches(j * 0.42),
                Inches(0.32),
                Inches(0.32),
            )
            _fill(cell, RGBColor(0x3A, 0x1A, 0x12) if blocked else NAVY_LIGHT)
            cell.line.color.rgb = CARD_BORDER
            cell.line.width = Pt(0.4)

    # Orange "failed GPS" trail into dead-end
    fail_path = [(1, 12), (2, 11), (3, 10), (4, 9), (5, 8), (6, 7), (7, 6), (8, 5), (8, 4)]
    for idx, (i, j) in enumerate(fail_path):
        d = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(7.7) + Inches(i * 0.42 + 0.07),
            Inches(0.9) + Inches(j * 0.42 + 0.07),
            Inches(0.18),
            Inches(0.18),
        )
        _fill(d, ORANGE if idx < len(fail_path) - 1 else RGBColor(0xEF, 0x44, 0x44))

    # X mark at dead-end
    xbox = _textbox(slide, Inches(10.9), Inches(2.4), Inches(0.8), Inches(0.5))
    _add_text(xbox.text_frame, "✕", size=22, bold=True, color=RGBColor(0xEF, 0x44, 0x44), align=PP_ALIGN.CENTER)

    badge = _round_rect(slide, Inches(7.85), Inches(0.35), Inches(4.6), Inches(0.4), CARD, ORANGE)
    btf = badge.text_frame
    btf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = btf.paragraphs[0].add_run()
    run.text = "SIMULATION VISUAL — ROAD NETWORK IN FLUX"
    _set_run(run, size=10, bold=True, color=ORANGE)

    # Left content
    _accent_bar(slide, MARGIN, Inches(0.55))
    eb = _textbox(slide, MARGIN, Inches(0.65), Inches(6.5), Inches(0.3))
    _add_text(eb.text_frame, "CASE STUDY", size=11, bold=True, color=ORANGE, space_after=0)

    title = _textbox(slide, MARGIN, Inches(1.0), Inches(6.6), Inches(1.0))
    _add_text(
        title.text_frame,
        "Operation Intramuros 2025",
        size=28,
        bold=True,
        color=WHITE,
        space_after=4,
    )
    _add_text(
        title.text_frame,
        "Simulation: M7.2 earthquake strikes the historic core of Manila.",
        size=14,
        color=SLATE,
        space_after=0,
        clear=False,
    )

    facts = [
        ("T+60s", "45% of roads blocked or impassable"),
        ("GPS Fail", "Standard navigation leads into dead-ends"),
        ("Cost", "Gridlock converts minutes into lives lost"),
    ]
    y = Inches(2.5)
    for k, v in facts:
        card = _round_rect(slide, MARGIN, y, Inches(6.5), Inches(1.05), CARD, CARD_BORDER)
        _rect(slide, MARGIN, y, Inches(0.1), Inches(1.05), ORANGE)
        kbox = _textbox(slide, MARGIN + Inches(0.3), y + Inches(0.18), Inches(6.0), Inches(0.3))
        _add_text(kbox.text_frame, k, size=12, bold=True, color=ORANGE, space_after=0)
        vbox = _textbox(slide, MARGIN + Inches(0.3), y + Inches(0.48), Inches(6.0), Inches(0.4))
        _add_text(vbox.text_frame, v, size=16, color=WHITE, space_after=0)
        y += Inches(1.2)

    _footer(slide, 3)


def slide_04_solution(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _abstract_flow(slide, intensity="subtle")
    _header(
        slide,
        "The Solution",
        "The Hybrid Advantage",
        subtitle="QuantumRelief Architecture — prediction meets precision pathfinding",
    )

    # Architecture flow: QML -> Fusion -> Dijkstra
    boxes = [
        ("QML LAYER", "Non-harmonic feature\nprediction of weight\nevolution under chaos", ORANGE),
        ("FUSION CORE", "Local state + predicted\nedge dynamics → next-\nhop decision policy", ACCENT_TEAL),
        ("CLASSICAL", "Dijkstra local\npathfinding on live,\npredicted weights", SLATE_LIGHT),
    ]
    bw = Inches(3.4)
    y = Inches(2.15)
    for i, (title, body, accent) in enumerate(boxes):
        x = MARGIN + i * (bw + Inches(0.55))
        card = _round_rect(slide, x, y, bw, Inches(2.6), CARD, CARD_BORDER)
        _rect(slide, x, y, bw, Inches(0.08), accent)
        t = _textbox(slide, x + Inches(0.25), y + Inches(0.35), bw - Inches(0.5), Inches(0.4))
        _add_text(t.text_frame, title, size=14, bold=True, color=accent, align=PP_ALIGN.CENTER)
        b = _textbox(slide, x + Inches(0.3), y + Inches(0.95), bw - Inches(0.6), Inches(1.4))
        tf = b.text_frame
        tf.word_wrap = True
        _add_text(tf, body, size=15, color=SLATE_LIGHT, align=PP_ALIGN.CENTER)

        if i < 2:
            arrow = _textbox(slide, x + bw - Inches(0.05), y + Inches(1.0), Inches(0.55), Inches(0.4))
            _add_text(arrow.text_frame, "→", size=24, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    # Gold standard callout
    call = _round_rect(slide, MARGIN, Inches(5.15), Inches(12.2), Inches(1.4), CARD, CARD_BORDER)
    _rect(slide, MARGIN, Inches(5.15), Inches(12.2), Inches(0.08), ORANGE)
    cbox = _textbox(slide, MARGIN + Inches(0.4), Inches(5.35), Inches(11.4), Inches(1.05))
    tf = cbox.text_frame
    tf.word_wrap = True
    _add_text(tf, "WHY HYBRID IS THE GOLD STANDARD", size=11, bold=True, color=ORANGE, space_after=6)
    _add_text(
        tf,
        "Quantum expressivity captures non-linear, non-harmonic disaster dynamics that classical "
        "models miss — while Dijkstra guarantees locally optimal traversal on the predicted graph. "
        "Neither alone is enough. Together, they navigate chaos before it hardens into gridlock.",
        size=14,
        color=SLATE_LIGHT,
        space_after=0,
        clear=False,
    )
    _footer(slide, 4)


def slide_05_hybrid_vs_dijkstra(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _abstract_flow(slide, intensity="subtle")
    _header(
        slide,
        "Comparative Logic",
        "Why Hybrid > Pure Dijkstra?",
        subtitle="Two operating models — one survives flux",
    )

    # Two columns
    cols = [
        (
            "PURE DIJKSTRA",
            SLATE,
            [
                "Requires near-global knowledge of the road graph",
                "Relies on static or lagging edge weights",
                "Recompute cost explodes as topology shifts",
                "Fails when the map is already obsolete",
            ],
            "Brittle in disaster flux",
        ),
        (
            "QUANTUMRELIEF",
            ORANGE,
            [
                "Operates on local neighborhood information",
                "Predicts weight evolution before it happens",
                "Sub-millisecond next-hop inference under pressure",
                "Routes around chaos that hasn't fully formed yet",
            ],
            "Anticipatory & resilient",
        ),
    ]
    cw = Inches(5.85)
    for i, (title, accent, bullets, tag) in enumerate(cols):
        x = MARGIN + i * (cw + Inches(0.45))
        card = _round_rect(slide, x, Inches(1.95), cw, Inches(4.55), CARD, CARD_BORDER)
        _rect(slide, x, Inches(1.95), cw, Inches(0.1), accent)

        t = _textbox(slide, x + Inches(0.35), Inches(2.25), cw - Inches(0.7), Inches(0.4))
        _add_text(t.text_frame, title, size=16, bold=True, color=accent, space_after=0)

        bbox = _textbox(slide, x + Inches(0.35), Inches(2.85), cw - Inches(0.7), Inches(2.8))
        _bullet_block(bbox.text_frame, bullets, size=14, color=SLATE_LIGHT, bullet_color=accent)

        tag_box = _round_rect(
            slide, x + Inches(0.35), Inches(5.75), cw - Inches(0.7), Inches(0.45), NAVY_MID, accent
        )
        tf = tag_box.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = tag.upper()
        _set_run(run, size=11, bold=True, color=accent)

    _footer(slide, 5)


def slide_06_tech_dive(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _abstract_flow(slide, intensity="subtle")
    _header(
        slide,
        "Technical Deep Dive",
        "Expressivity & Efficiency",
        subtitle="Built for high-dimensional disaster features — and real-time decisions",
    )

    tiles = [
        ("◎", "FiLM Conditioning", "Epicenter-aware Feature-wise Linear Modulation (γ, β) adapts the network to disaster geometry in real time."),
        ("⟳", "Data Re-uploading", "Repeated encoding maps high-dimensional edge features into expressive quantum Hilbert space."),
        ("Σ", "Quantum Contribution", "Live % = 100×mean(|W_q|)/(mean(|W_c|)+mean(|W_q|)) from PHN combine — ≈37.9% on trained film_hybrid.pt."),
        ("▣", "36-D Local State", "Table I feature layout: epicenter, start/dest, and up to 5 adjacent edges with topology cues."),
        ("⬡", "Hybrid PHN", "PennyLane quantum scaffold fused with classical FiLM for production-ready fallback."),
        ("⏱", "Latency & QPU", "Hybrid slower on classical simulators; roadmap: real QPU accelerates routing operators."),
    ]

    cols, rows = 3, 2
    tw, th = Inches(3.85), Inches(2.15)
    gap_x, gap_y = Inches(0.3), Inches(0.25)
    origin_y = Inches(1.95)
    for idx, (icon, title, body) in enumerate(tiles):
        r, c = divmod(idx, cols)
        x = MARGIN + c * (tw + gap_x)
        y = origin_y + r * (th + gap_y)
        card = _round_rect(slide, x, y, tw, th, CARD, CARD_BORDER)
        _rect(slide, x, y, Inches(0.08), th, ORANGE)

        ib = _textbox(slide, x + Inches(0.25), y + Inches(0.2), Inches(0.5), Inches(0.4))
        _add_text(ib.text_frame, icon, size=18, bold=True, color=ORANGE, space_after=0)

        tb = _textbox(slide, x + Inches(0.7), y + Inches(0.25), tw - Inches(0.95), Inches(0.35))
        _add_text(tb.text_frame, title, size=14, bold=True, color=WHITE, space_after=0)

        bb = _textbox(slide, x + Inches(0.25), y + Inches(0.75), tw - Inches(0.5), Inches(1.2))
        tf = bb.text_frame
        tf.word_wrap = True
        _add_text(tf, body, size=12, color=SLATE, space_after=0)

    _footer(slide, 6)


def slide_07_demo(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Full dark background "cinema" feel
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0x06, 0x10, 0x18))

    # Subtle grid overlay across full slide
    for i in range(0, 28):
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(i * 0.5), 0, Inches(0.01), SLIDE_H
        )
        _fill(line, RGBColor(0x12, 0x28, 0x3A))
    for j in range(0, 16):
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, Inches(j * 0.5), SLIDE_W, Inches(0.01)
        )
        _fill(line, RGBColor(0x12, 0x28, 0x3A))

    # Vignette panels
    _rect(slide, 0, 0, SLIDE_W, Inches(1.1), RGBColor(0x04, 0x0C, 0x14))
    _rect(slide, 0, Inches(6.4), SLIDE_W, Inches(1.1), RGBColor(0x04, 0x0C, 0x14))

    # Simulated route trails
    for k, color in enumerate([ORANGE, ACCENT_TEAL, ORANGE_SOFT]):
        for n in range(18):
            d = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(1.2 + n * 0.55 + k * 0.08),
                Inches(2.8 + 0.35 * ((n + k * 2) % 5) + k * 0.55),
                Inches(0.14),
                Inches(0.14),
            )
            _fill(d, color)

    # Center play button motif
    ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.55), Inches(2.55), Inches(2.25), Inches(2.25))
    _fill(ring, CARD)
    _line(ring, ORANGE, 2.5)

    inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.85), Inches(2.85), Inches(1.65), Inches(1.65))
    _fill(inner, NAVY)
    _line(inner, ORANGE, 1.5)

    play = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(6.35), Inches(3.25), Inches(0.85), Inches(0.85)
    )
    _fill(play, ORANGE)
    play.rotation = 90.0

    # Captions
    top = _textbox(slide, Inches(0.8), Inches(0.35), Inches(11.5), Inches(0.5))
    _add_text(
        top.text_frame,
        "DEMO VIDEO  ·  REAL-TIME MANILA EVACUATION",
        size=14,
        bold=True,
        color=ORANGE,
        align=PP_ALIGN.CENTER,
        space_after=0,
    )

    mid = _textbox(slide, Inches(2.0), Inches(5.05), Inches(9.3), Inches(0.7))
    _add_text(
        mid.text_frame,
        "Simulation of vehicles navigating Manila — avoiding dynamic blockages in real time.",
        size=16,
        color=SLATE_LIGHT,
        align=PP_ALIGN.CENTER,
        space_after=0,
    )

    bottom = _textbox(slide, Inches(2.0), Inches(6.65), Inches(9.3), Inches(0.4))
    _add_text(
        bottom.text_frame,
        "▶  PLACEHOLDER — Insert live demo / screen recording here",
        size=12,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        space_after=0,
    )


def slide_08_results(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _abstract_flow(slide, intensity="subtle")
    _header(
        slide,
        "The Results",
        "Proving the Edge",
        subtitle="Hybrid QuantumRelief vs. Pure Classical baselines",
    )

    # Chart
    chart_data = CategoryChartData()
    chart_data.categories = ["Accuracy", "Arrival Rate"]
    chart_data.add_series("Hybrid (QuantumRelief)", (94, 95))
    chart_data.add_series("Pure Classical", (87, 92))

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.55),
        Inches(1.9),
        Inches(7.6),
        Inches(4.6),
        chart_data,
    )
    chart = chart_frame.chart
    _style_chart(chart, value_axis_max=100)

    # Right callouts
    metrics = [
        ("94%", "Hybrid Accuracy", "vs 87% classical"),
        ("95%", "Arrival Rate", "vs 92% classical"),
    ]
    y = Inches(1.95)
    for num, title, sub in metrics:
        card = _round_rect(slide, Inches(8.5), y, Inches(4.25), Inches(1.45), CARD, CARD_BORDER)
        _rect(slide, Inches(8.5), y, Inches(0.1), Inches(1.45), ORANGE)
        n = _textbox(slide, Inches(8.8), y + Inches(0.25), Inches(3.7), Inches(0.55))
        _add_text(n.text_frame, num, size=32, bold=True, color=ORANGE, space_after=0)
        t = _textbox(slide, Inches(8.8), y + Inches(0.8), Inches(3.7), Inches(0.5))
        _add_text(t.text_frame, f"{title}  ·  {sub}", size=12, color=SLATE_LIGHT, space_after=0)
        y += Inches(1.65)

    # Summary callout
    call = _round_rect(slide, Inches(8.5), Inches(5.2), Inches(4.25), Inches(1.35), CARD, ORANGE)
    cbox = _textbox(slide, Inches(8.75), Inches(5.4), Inches(3.8), Inches(1.05))
    tf = cbox.text_frame
    tf.word_wrap = True
    _add_text(tf, "KEY INSIGHT", size=10, bold=True, color=ORANGE, space_after=4)
    _add_text(
        tf,
        "Quantum neurons contribute 45.3% to decision accuracy.",
        size=14,
        bold=True,
        color=WHITE,
        space_after=0,
        clear=False,
    )
    _footer(slide, 8)


def slide_09_monetization(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _abstract_flow(slide, intensity="subtle")
    _header(
        slide,
        "Strategic Monetization",
        "Scaling Resilience",
        subtitle="Dual go-to-market: B2G civic infrastructure + B2B logistics platforms",
    )

    # Two segment headers
    for i, (seg, label) in enumerate([("B2G", "GOVERNMENT"), ("B2B", "LOGISTICS")]):
        x = MARGIN + i * Inches(6.25)
        pill = _round_rect(slide, x, Inches(1.9), Inches(2.0), Inches(0.38), CARD, ORANGE)
        tf = pill.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = f"{seg}  ·  {label}"
        _set_run(run, size=11, bold=True, color=ORANGE)

    products = [
        (
            "Simulation-as-a-Service",
            "B2G",
            "City planners stress-test evacuation plans against synthetic M7+ scenarios before the real event.",
        ),
        (
            "Quantum Routing API",
            "B2B",
            "Plug-in resilience layer for Grab / logistics fleets — dynamic rerouting when the city breaks.",
        ),
        (
            "Insurance Risk Modeling",
            "B2G / B2B",
            "Quantify corridor failure probability and expected delay loss for underwriters and agencies.",
        ),
    ]
    y = Inches(2.5)
    for title, tag, body in products:
        card = _round_rect(slide, MARGIN, y, Inches(12.2), Inches(1.2), CARD, CARD_BORDER)
        _rect(slide, MARGIN, y, Inches(0.1), Inches(1.2), ORANGE)
        t = _textbox(slide, MARGIN + Inches(0.35), y + Inches(0.2), Inches(8.5), Inches(0.35))
        _add_text(t.text_frame, title, size=16, bold=True, color=WHITE, space_after=0)
        tag_b = _textbox(slide, Inches(10.0), y + Inches(0.22), Inches(2.4), Inches(0.3))
        _add_text(tag_b.text_frame, tag, size=11, bold=True, color=ORANGE, align=PP_ALIGN.RIGHT, space_after=0)
        b = _textbox(slide, MARGIN + Inches(0.35), y + Inches(0.6), Inches(11.4), Inches(0.45))
        _add_text(b.text_frame, body, size=13, color=SLATE, space_after=0)
        y += Inches(1.35)

    _footer(slide, 9)


def slide_10_conclusion(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _abstract_flow(slide, intensity="strong")
    _rect(slide, 0, 0, Inches(0.12), SLIDE_H, ORANGE)

    eyebrow = _textbox(slide, MARGIN, Inches(1.6), Inches(12), Inches(0.35))
    _add_text(
        eyebrow.text_frame,
        "CONCLUSION",
        size=12,
        bold=True,
        color=ORANGE,
        align=PP_ALIGN.CENTER,
        space_after=0,
    )

    tag = _textbox(slide, MARGIN, Inches(2.2), Inches(12.2), Inches(1.2))
    _add_text(
        tag.text_frame,
        "Quantum Intelligence. Human Relief.",
        size=36,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        space_after=0,
    )

    _accent_bar(slide, Inches(5.7), Inches(3.5), width=Inches(1.9), height=Inches(0.06))

    contact = _textbox(slide, MARGIN, Inches(3.9), Inches(12.2), Inches(0.9))
    tf = contact.text_frame
    _add_text(
        tf,
        "contact@quantumrelief.ai  ·  www.quantumrelief.ai",
        size=16,
        color=SLATE_LIGHT,
        align=PP_ALIGN.CENTER,
        space_after=10,
    )
    _add_text(
        tf,
        "Team 5 — Quantrio  |  QC4SG — SEA Hackathon",
        size=14,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        space_after=0,
        clear=False,
    )

    # Bottom brand strip
    strip = _round_rect(slide, Inches(3.8), Inches(5.5), Inches(5.7), Inches(0.7), CARD, CARD_BORDER)
    stf = strip.text_frame
    stf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = stf.paragraphs[0].add_run()
    run.text = "QuantumRelief"
    _set_run(run, size=20, bold=True, color=ORANGE)


def build(output: Path = OUT_PATH) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_case_study(prs)
    slide_04_solution(prs)
    slide_05_hybrid_vs_dijkstra(prs)
    slide_06_tech_dive(prs)
    slide_07_demo(prs)
    slide_08_results(prs)
    slide_09_monetization(prs)
    slide_10_conclusion(prs)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return output


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path} ({path.stat().st_size} bytes)")
